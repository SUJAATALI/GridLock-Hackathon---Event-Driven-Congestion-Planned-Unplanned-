"""
Generates a boilerplate pitch deck (.pptx) for the Gridlock 2.0 submission.
Dark control-room theme, real verified numbers, tables, and speaker notes.
The teammate edits this in PowerPoint/Keynote/Google Slides — it's a starting point,
not a finished deck (placeholders for screenshots are marked).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (matches the dashboard) ----
BG     = RGBColor(0x0D, 0x11, 0x17)
PANEL  = RGBColor(0x16, 0x1B, 0x22)
LINE   = RGBColor(0x30, 0x36, 0x3D)
INK    = RGBColor(0xE6, 0xED, 0xF3)
MUTED  = RGBColor(0x8B, 0x94, 0x9E)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
HI     = RGBColor(0xF8, 0x51, 0x49)
MED    = RGBColor(0xD2, 0x99, 0x22)
LOW    = RGBColor(0x3F, 0xB9, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "Arial"
    return tb


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def chip(s, x, y, label, color):
    c = box(s, x, y, 0.18, 0.18, fill=color)
    c.adjustments[0] = 0.5  # rounded -> circle-ish (rounded rect already type 1; ignore)
    return c


def table(s, x, y, w, rows, col_w, header_fill=PANEL, row_h=0.5, fontsize=16):
    """rows: list of rows; row = list of (text, color, bold). First row = header."""
    n = len(rows)
    gt = s.shapes.add_table(n, len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(row_h * n)).table
    # kill default styling
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        gt.rows[ri].height = Inches(row_h)
        for ci, (val, color, bold) in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri == 0 else BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(fontsize); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Arial"
    return gt


def kicker(s, txt):
    text(s, 0.7, 0.45, 11, 0.4, [[(txt, 13, ACCENT, True)]])


# ============================================================ SLIDE 1 — Title
s = slide()
text(s, 0.7, 2.4, 12, 1.4, [[("🚦 Event Congestion Command Center", 44, INK, True)]])
text(s, 0.7, 3.7, 12, 0.8,
     [[("Predicting traffic-event impact for the Bengaluru Traffic Police", 22, MUTED, False)]])
box(s, 0.72, 4.6, 4.2, 0.04, fill=ACCENT)
text(s, 0.7, 4.8, 12, 1.2, [
    [("Flipkart Gridlock 2.0  ×  BTP", 16, INK, True)],
    [("Built on real ASTraM data — 8,173 actual traffic events (Nov 2023 – Apr 2024). No simulations.", 14, MUTED, False)],
])
notes(s, "Bengaluru loses enormous time to traffic events — breakdowns, tree falls, construction, "
         "rallies. The police need to know WHICH events will shut a road and WHERE to put their people. "
         "We built a command center on the REAL provided dataset — 8,173 events — not simulated data.")

# ============================================================ SLIDE 2 — Data & enrichment
s = slide()
kicker(s, "THE DATA")
text(s, 0.7, 0.85, 12, 0.8, [[("We used only the provided data — then enriched it", 30, INK, True)]])
text(s, 0.7, 1.7, 12, 0.5,
     [[("Raw ASTraM had ~46 columns, many 98–99% empty. We engineered what's real — via 3 reproducible scripts.",
        15, MUTED, False)]])
# pipeline strip
box(s, 0.9, 2.5, 2.4, 0.9, fill=PANEL, line=LINE)
text(s, 0.9, 2.5, 2.4, 0.9, [[("RAW", 18, INK, True)], [("46 columns", 13, MUTED, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 3.4, 2.5, 0.7, 0.9, [[("→", 28, ACCENT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 4.1, 2.5, 4.6, 0.9, fill=PANEL, line=ACCENT)
text(s, 4.1, 2.5, 4.6, 0.9, [[("3 enrichment scripts", 18, ACCENT, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 8.7, 2.5, 0.7, 0.9, [[("→", 28, ACCENT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 9.4, 2.5, 3.0, 0.9, fill=PANEL, line=LINE)
text(s, 9.4, 2.5, 3.0, 0.9, [[("ENRICHED", 18, LOW, True)], [("55 columns · 8,173 rows", 12, MUTED, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# three script cards
cards = [
    ("add_locality_hotspot.py", "locality_hotspot", "Cracked the 38% \"Non-corridor\" blob into real neighborhoods — Jayanagar, KR Puram, Whitefield (1 km GPS grid + address)."),
    ("add_region.py", "region", "8 city zones from a 3×3 grid over the data's own lat/long — central zone holds ~50% of events."),
    ("add_time_features.py", "hour · day · weekend · night · time_band", "Exposed the rhythm: 2 AM freight peak + 10 AM–noon rush peak."),
]
for i, (fn, cols, desc) in enumerate(cards):
    x = 0.9 + i * 4.05
    box(s, x, 3.8, 3.8, 2.7, fill=PANEL, line=LINE)
    text(s, x + 0.25, 4.0, 3.3, 2.4, [
        [(fn, 14, ACCENT, True)],
        [("", 6, MUTED, False)],
        [(cols, 13, LOW, True)],
        [("", 6, MUTED, False)],
        [(desc, 13, INK, False)],
    ], space=1.05)
text(s, 0.7, 6.75, 12, 0.5,
     [[("We enriched the data, we didn't replace it — re-runnable from the raw file.", 14, MUTED, True)]])
notes(s, "We used ONLY the provided ASTraM dataset. Raw had ~46 columns but many were 98-99% empty "
         "(resolution time, assigned officer) so we were honest about what's usable and enriched it. "
         "Three documented scripts — locality names, city zones, and time patterns — all derived "
         "purely from the provided dataset (lat/long + timestamp). Fully reproducible from the raw "
         "file, offline. The single strongest predictor turned out to be the event's CAUSE.")

# ============================================================ SLIDE 3 — Bake-off
s = slide()
kicker(s, "MODEL SELECTION")
text(s, 0.7, 0.85, 12, 0.8, [[("We tested 3 models — a fair bake-off", 30, INK, True)]])
text(s, 0.7, 1.7, 12, 0.5,
     [[("Same data, same 75/25 split, same imbalance handling. We let the data pick the winner.", 15, MUTED, False)]])
rows = [
    [("Model", INK, True), ("ROC-AUC", MUTED, True), ("PR-AUC  ⭐", ACCENT, True)],
    [("XGBoost   🏆", LOW, True), ("0.80", INK, True), ("0.44", LOW, True)],
    [("HistGradientBoosting", INK, False), ("0.77", MUTED, False), ("0.41", MUTED, False)],
    [("LightGBM", INK, False), ("0.78", MUTED, False), ("0.41", MUTED, False)],
]
table(s, 0.9, 2.6, 7.5, rows, [3.7, 1.9, 1.9], row_h=0.62, fontsize=17)
box(s, 8.9, 2.6, 3.5, 2.48, fill=PANEL, line=LINE)
text(s, 9.15, 2.8, 3.0, 2.1, [
    [("Why PR-AUC?", 16, ACCENT, True)],
    [("", 6, MUTED, False)],
    [("Closures are rare (8.3%).", 14, INK, False)],
    [("PR-AUC is the honest metric", 14, INK, False)],
    [("for catching rare events —", 14, INK, False)],
    [("so it's our tiebreaker.", 14, INK, False)],
], space=1.05)
text(s, 0.9, 5.5, 11.5, 0.8,
     [[("XGBoost won on PR-AUC — the metric that matters for rare events. That's the model we kept.",
        18, INK, True)]])
notes(s, "We didn't pick a model and hope. Three industry-standard gradient-boosted models, identical "
         "conditions. XGBoost won on PR-AUC — the metric that matters when the thing you're predicting "
         "is rare. So that's our closure-risk engine.")

# ============================================================ SLIDE 4 — Performance honest
s = slide()
kicker(s, "MODEL PERFORMANCE")
text(s, 0.7, 0.85, 12, 0.8, [[("How good is it — measured honestly", 30, INK, True)]])
box(s, 0.9, 1.75, 11.5, 0.95, fill=PANEL, line=MED)
text(s, 1.15, 1.85, 11, 0.8, [
    [("⚠  Only 8.3% of events cause a closure. ", 15, MED, True),
     ("So \"accuracy\" is misleading — a model that says \"never close\" scores 92% and is useless.", 15, INK, False)],
    [("We measure what counts: do we catch the real closures?", 15, INK, True)],
], space=1.0)
rows = [
    [("Metric", INK, True), ("Score", ACCENT, True), ("What it means", MUTED, True)],
    [("PR-AUC", INK, True), ("0.44", LOW, True), ("~5× better than random (0.08)", INK, False)],
    [("ROC-AUC", INK, True), ("0.80", LOW, True), ("Ranks a closure above a non-closure 4 of 5 times", INK, False)],
    [("Recall", INK, True), ("54%", LOW, True), ("Catches over half of all real closures", INK, False)],
    [("Backtest", INK, True), ("4.2×", LOW, True), ("Flagged events actually closed 4.2× more often", INK, False)],
]
table(s, 0.9, 3.0, 11.5, rows, [2.2, 1.6, 7.7], row_h=0.62, fontsize=16)
text(s, 0.7, 6.5, 12, 0.6,
     [[("“We'd rather show a true 0.44 than a fake 0.99.” Every number is on data the model never saw.",
        17, ACCENT, True)]])
notes(s, "Because closures are rare, accuracy is a trap — 92% is the do-nothing baseline. We report "
         "PR-AUC and recall instead: ~5x better than chance, catches 54% of closures, and history shows "
         "flagged events closed 4.2x more often. Honest, tested, on unseen data.")

# ============================================================ SLIDE 5 — 3-model system
s = slide()
kicker(s, "THE SYSTEM")
text(s, 0.7, 0.85, 12, 0.8, [[("Three models, one command center", 30, INK, True)]])
mods = [
    ("MODEL 1", "Closure Risk", "Will this event shut the road?", "XGBoost · PR-AUC 0.44, catches 54% of closures", ACCENT),
    ("MODEL 2", "Hotspot Forecaster", "Where & when do events cluster?", "By neighborhood × hour × day · beats naive avg 27–38%", LOW),
    ("MODEL 3", "Recommendation Engine", "What should the police do?", "Officers · barricade · pre-position · backtested 4.2× lift", MED),
]
for i, (tag, title, q, detail, col) in enumerate(mods):
    x = 0.9 + i * 4.05
    box(s, x, 2.2, 3.8, 3.4, fill=PANEL, line=LINE)
    box(s, x, 2.2, 3.8, 0.12, fill=col)
    text(s, x + 0.3, 2.55, 3.2, 3.0, [
        [(tag, 13, col, True)],
        [("", 4, MUTED, False)],
        [(title, 20, INK, True)],
        [("", 8, MUTED, False)],
        [(q, 15, MUTED, True)],
        [("", 8, MUTED, False)],
        [(detail, 13, INK, False)],
    ], space=1.05)
    if i < 2:
        text(s, x + 3.8, 2.2, 0.25, 3.4, [[("→", 22, ACCENT, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 6.2, 12, 0.5,
     [[("Predict impact  →  locate hotspots  →  recommend response.", 16, MUTED, True)]])
notes(s, "Three models working together. Model 1 scores an incoming event's closure risk. Model 2 tells "
         "us where and when events cluster, by real neighborhood. Model 3 turns risk into action: how many "
         "officers, barricade or not, and where to pre-position. All explainable — no black box.")

# ============================================================ SLIDE 6 — Dashboard
s = slide()
kicker(s, "LIVE DEMO")
text(s, 0.7, 0.85, 12, 0.8, [[("The dashboard — report an event, get a plan", 30, INK, True)]])
box(s, 0.9, 1.9, 11.5, 4.4, fill=PANEL, line=ACCENT)
text(s, 0.9, 3.7, 11.5, 1.0, [
    [("[  SCREENSHOT: dashboard with map + Report Event form + recommendation panel  ]", 16, MUTED, True)],
    [("replace with a real screenshot from http://127.0.0.1:5050", 12, MUTED, False)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 6.5, 12, 0.5,
     [[("Report an event → closure risk % + officer/barricade recommendation + city hotspot map, in real time.",
        15, INK, True)]])
notes(s, "LIVE DEMO. Open the dashboard, hit 'Replay sample events' for a hands-free walkthrough — events "
         "stream onto the map spanning low to high risk. Show: report form -> risk % -> recommendation panel "
         "-> hotspot heatmap (drag the hour slider to reveal the 2 AM & 10-12 peaks) -> pre-positioning list. "
         "Use the replay button, not live typing, so it's deterministic.")

# ============================================================ SLIDE 7 — Honesty
s = slide()
kicker(s, "WHY YOU CAN TRUST THIS")
text(s, 0.7, 0.85, 12, 0.8, [[("Honesty is our edge", 30, INK, True)]])
pts = [
    ("🎯", "Real held-out metrics", "Every score is on 25% of data the model never saw. We report PR-AUC & recall — right for rare events — not flattering accuracy."),
    ("🐛", "We caught & fixed a leakage bug", "A timestamp-parsing glitch was secretly inflating our score (116 events). We found it, fixed it at the source, and report the honest post-fix number."),
    ("📋", "Recommendations are backtested, not “optimal”", "No deployment ground truth exists in the data, so we never claim optimal. We verified vs history: flagged events closed 4.2× more often."),
    ("📦", "Only the provided data", "Every feature is engineered purely from the ASTraM dataset — location, time, cause. No external datasets. The real driver is event cause."),
]
for i, (ic, title, desc) in enumerate(pts):
    x = 0.9 + (i % 2) * 5.9
    y = 2.0 + (i // 2) * 2.1
    box(s, x, y, 5.6, 1.85, fill=PANEL, line=LINE)
    text(s, x + 0.25, y + 0.18, 5.1, 1.6, [
        [(ic + "  ", 18, INK, False), (title, 16, ACCENT, True)],
        [("", 4, MUTED, False)],
        [(desc, 13, INK, False)],
    ], space=1.05)
text(s, 0.7, 6.55, 12, 0.6,
     [[("“We'd rather present a true 0.44 than a fake 0.99.”", 18, ACCENT, True)]], align=PP_ALIGN.CENTER)
notes(s, "The slide we're proudest of. We could have shown 92% accuracy and hoped. Instead: honest metrics, "
         "we caught our own leakage bug, recommendations backtested against real history, and every feature "
         "comes only from the provided data. If you're judging on rigor, this is where we earn it.")

# ============================================================ SLIDE 8 — Impact + roadmap
s = slide()
kicker(s, "IMPACT & ROADMAP")
text(s, 0.7, 0.85, 12, 0.8, [[("What it means for BTP — and what's next", 30, INK, True)]])
box(s, 0.9, 2.0, 5.6, 3.6, fill=PANEL, line=LINE)
text(s, 1.15, 2.25, 5.1, 3.2, [
    [("IMPACT TODAY", 15, LOW, True)],
    [("", 6, MUTED, False)],
    [("• Proactive: pre-position patrols at predicted hotspots before incidents happen.", 14, INK, False)],
    [("", 4, MUTED, False)],
    [("• Reactive: instant closure risk + response the moment an event is reported.", 14, INK, False)],
    [("", 4, MUTED, False)],
    [("• Explainable — officers can see why, so they trust it.", 14, INK, False)],
], space=1.1)
box(s, 6.8, 2.0, 5.6, 3.6, fill=PANEL, line=LINE)
text(s, 7.05, 2.25, 5.1, 3.2, [
    [("ROADMAP  (not yet built)", 15, MED, True)],
    [("", 6, MUTED, False)],
    [("• Live ASTraM feed integration", 14, INK, False)],
    [("", 4, MUTED, False)],
    [("• MapmyIndia for map, traffic & diversion routing", 14, INK, False)],
    [("", 4, MUTED, False)],
    [("• \"Blast radius\" impact-zone visualization", 14, INK, False)],
    [("", 4, MUTED, False)],
    [("• Learning loop from real deployment outcomes", 14, INK, False)],
], space=1.1)
text(s, 0.7, 6.5, 12, 0.5,
     [[("Honest about what's live today vs. what comes next.", 14, MUTED, True)]])
notes(s, "Impact: proactive + reactive, and explainable so officers trust it. Roadmap is clearly labeled "
         "NOT YET BUILT — live feed, MapmyIndia routing/diversions, blast-radius zone, and a learning loop. "
         "We keep the line honest between what works today and what's future.")

OUT = "Gridlock2.0_CommandCenter_DECK.pptx"
prs.save(OUT)
print("saved", OUT, "—", len(prs.slides._sldIdLst), "slides")
